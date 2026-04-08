"""
LMS Integration — Connexion Moodle (+ Canvas futur)
Permet aux étudiants de connecter leur plateforme universitaire
et d'importer le contenu de leurs cours pour alimenter les quiz/flashcards/résumés.
"""

import re
import logging
import requests
from io import BytesIO

log = logging.getLogger(__name__)

MOODLE_SERVICE = 'moodle_mobile_app'

# ─── Universités belges connues ───────────────────────────────────────────────
KNOWN_UNIVERSITIES = [
    {'name': 'UCLouvain', 'url': 'https://moodleucl.uclouvain.be', 'platform': 'moodle'},
    {'name': 'ULB', 'url': 'https://uv.ulb.ac.be', 'platform': 'moodle'},
    {'name': 'ULiège', 'url': 'https://lola.uliege.be', 'platform': 'moodle'},
    {'name': 'UNamur', 'url': 'https://webcampus.unamur.be', 'platform': 'moodle'},
    {'name': 'USaint-Louis', 'url': 'https://moodle.usaintlouis.be', 'platform': 'moodle'},
    {'name': 'UMons', 'url': 'https://moodle.umons.ac.be', 'platform': 'moodle'},
    {'name': 'HE Vinci', 'url': 'https://moodle.vinci.be', 'platform': 'moodle'},
    {'name': 'EPHEC', 'url': 'https://moodle.ephec.be', 'platform': 'moodle'},
    {'name': 'ICHEC', 'url': 'https://moodle.ichec.be', 'platform': 'moodle'},
    {'name': 'KU Leuven', 'url': 'https://toledo.kuleuven.be', 'platform': 'toledo'},
    {'name': 'VUB', 'url': 'https://canvas.vub.be', 'platform': 'canvas'},
    {'name': 'UGent', 'url': 'https://ufora.ugent.be', 'platform': 'ufora'},
    {'name': 'UAntwerpen', 'url': 'https://blackboard.uantwerpen.be', 'platform': 'blackboard'},
]


# ─── Moodle Authentication ───────────────────────────────────────────────────

def moodle_authenticate(site_url: str, username: str, password: str) -> str:
    """Authenticate with Moodle and return a web service token."""
    url = f"{site_url.rstrip('/')}/login/token.php"
    try:
        r = requests.post(url, data={
            'username': username,
            'password': password,
            'service': MOODLE_SERVICE,
        }, timeout=15)
        data = r.json()
    except requests.RequestException as e:
        raise ValueError(f"Impossible de joindre {site_url}: {e}")

    if 'token' in data:
        return data['token']
    error = data.get('error', data.get('message', 'Échec de connexion'))
    raise ValueError(f"Moodle: {error}")


def moodle_call(site_url: str, token: str, function: str, **params) -> dict:
    """Call a Moodle Web Service function."""
    url = f"{site_url.rstrip('/')}/webservice/rest/server.php"
    payload = {
        'wstoken': token,
        'wsfunction': function,
        'moodlewsrestformat': 'json',
        **params,
    }
    try:
        r = requests.post(url, data=payload, timeout=20)
        data = r.json()
    except requests.RequestException as e:
        raise ValueError(f"Erreur API Moodle: {e}")

    if isinstance(data, dict) and data.get('exception'):
        raise ValueError(data.get('message', 'Erreur Moodle'))
    return data


# ─── Moodle Data Fetching ─────────────────────────────────────────────────────

def get_site_info(site_url: str, token: str) -> dict:
    """Get Moodle site info (name, user fullname, etc.)."""
    data = moodle_call(site_url, token, 'core_webservice_get_site_info')
    return {
        'site_name': data.get('sitename', ''),
        'user_fullname': data.get('fullname', ''),
        'moodle_user_id': data.get('userid'),
    }


def get_courses(site_url: str, token: str, moodle_user_id: int) -> list:
    """Get enrolled courses for the user."""
    courses = moodle_call(site_url, token, 'core_enrol_get_users_courses', userid=moodle_user_id)
    if not isinstance(courses, list):
        return []
    return [{
        'id': c['id'],
        'name': c.get('fullname', ''),
        'shortname': c.get('shortname', ''),
        'category': c.get('categoryname', ''),
    } for c in courses]


def get_course_content(site_url: str, token: str, course_id: int) -> list:
    """Get course content: sections → modules → files."""
    sections = moodle_call(site_url, token, 'core_course_get_contents', courseid=course_id)
    if not isinstance(sections, list):
        return []
    result = []
    for section in sections:
        sec = {
            'id': section.get('id'),
            'name': section.get('name', ''),
            'summary': _strip_html(section.get('summary', '')),
            'modules': [],
        }
        for mod in section.get('modules', []):
            modname = mod.get('modname', '')
            # Garder uniquement les modules utiles (resource, page, url, folder, book)
            if modname not in ('resource', 'page', 'url', 'folder', 'book', 'label'):
                continue
            module = {
                'id': mod.get('id'),
                'name': mod.get('name', ''),
                'type': modname,
                'description': _strip_html(mod.get('description', '')),
                'contents': [],
            }
            for content in mod.get('contents', []):
                filename = content.get('filename', '')
                mimetype = content.get('mimetype', '')
                # Filtrer : PDF, DOCX, PPTX, HTML, TXT
                if _is_importable(filename, mimetype):
                    module['contents'].append({
                        'filename': filename,
                        'fileurl': content.get('fileurl', ''),
                        'filesize': content.get('filesize', 0),
                        'mimetype': mimetype,
                    })
            # Garder le module s'il a du contenu ou une description
            if module['contents'] or module['description']:
                sec['modules'].append(module)
        if sec['modules'] or sec['summary']:
            result.append(sec)
    return result


def download_and_extract(site_url: str, token: str, file_url: str) -> str:
    """Download a file from Moodle and extract text content."""
    download_url = f"{file_url}{'&' if '?' in file_url else '?'}token={token}"

    try:
        r = requests.get(download_url, timeout=30, allow_redirects=True)
        r.raise_for_status()
    except requests.RequestException as e:
        raise ValueError(f"Impossible de télécharger: {e}")

    content_type = r.headers.get('content-type', '')
    raw = r.content

    # PDF
    if 'pdf' in content_type or file_url.endswith('.pdf'):
        return _extract_pdf(raw)

    # HTML / Text
    if 'html' in content_type or 'text' in content_type:
        text = raw.decode('utf-8', errors='ignore')
        return _strip_html(text)[:20000]

    # DOCX
    if 'word' in content_type or file_url.endswith('.docx'):
        return _extract_docx(raw)

    # PPTX
    if 'presentation' in content_type or file_url.endswith('.pptx'):
        return _extract_pptx(raw)

    return f'[Format non supporté: {content_type}]'


# ─── Text Extraction Helpers ─────────────────────────────────────────────────

def _strip_html(html: str) -> str:
    """Strip HTML tags and clean whitespace."""
    if not html:
        return ''
    text = re.sub(r'<[^>]+>', ' ', html)
    text = re.sub(r'\s+', ' ', text).strip()
    return text


def _is_importable(filename: str, mimetype: str) -> bool:
    """Check if a file is importable (PDF, DOCX, PPTX, HTML, TXT)."""
    ext = filename.rsplit('.', 1)[-1].lower() if '.' in filename else ''
    importable_exts = {'pdf', 'docx', 'doc', 'pptx', 'ppt', 'txt', 'html', 'htm', 'md'}
    importable_mimes = {'application/pdf', 'text/', 'application/vnd.openxmlformats'}
    return ext in importable_exts or any(m in mimetype for m in importable_mimes)


def _extract_pdf(raw: bytes) -> str:
    """Extract text from PDF bytes."""
    try:
        import fitz  # PyMuPDF
        doc = fitz.open(stream=raw, filetype='pdf')
        text = '\n'.join(page.get_text() for page in doc)
        doc.close()
        return text[:20000]
    except ImportError:
        # Fallback: essayer pdfminer
        try:
            from pdfminer.high_level import extract_text
            return extract_text(BytesIO(raw))[:20000]
        except ImportError:
            return '[PDF — extraction non disponible sur ce serveur]'


def _extract_docx(raw: bytes) -> str:
    """Extract text from DOCX bytes."""
    try:
        from docx import Document
        doc = Document(BytesIO(raw))
        text = '\n'.join(p.text for p in doc.paragraphs if p.text.strip())
        return text[:20000]
    except ImportError:
        return '[DOCX — python-docx non installé]'


def _extract_pptx(raw: bytes) -> str:
    """Extract text from PPTX bytes."""
    try:
        from pptx import Presentation
        prs = Presentation(BytesIO(raw))
        text = '\n'.join(
            shape.text for slide in prs.slides
            for shape in slide.shapes if hasattr(shape, 'text') and shape.text.strip()
        )
        return text[:20000]
    except ImportError:
        return '[PPTX — python-pptx non installé]'
